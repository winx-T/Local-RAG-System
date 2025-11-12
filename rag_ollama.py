import os
import re
import json
import hashlib
import csv
import logging
import sys
import signal
import shutil
import threading
from typing import List, Dict, Tuple, Optional, Any
from datetime import datetime, timedelta
from pathlib import Path
import pickle
import time
import traceback
import json as json_lib

# 🛠️ CONFIGURATION
class MyConfig:
    DEFAULTS = {
        'data_folder': './data',
        'cache_folder': './cache',
        'backup_folder': './backups',
        'log_folder': './logs',
        'ollama_url': 'http://localhost:11434',
        'model_name': 'llama3.2:3b',
        'embedding_model': 'all-MiniLM-L6-v2',
        'use_reranking': True,
        'chunk_size': 500,
        'chunk_overlap': 75,
        'top_k_retrieval': 20,
        'top_k_rerank': 6,
        'batch_size': 24,
        'cache_ttl_hours': 48,
        'max_cache_size_mb': 150,
        'ollama_timeout': 90,
        'max_file_size_mb': 100,
        'max_documents': 10000,
        'enable_cache': True,
        'show_sources': True,
        'stream_answers': True,
        'save_history': True,
        'auto_backup': True,
        'language': 'auto'
    }

    def __init__(self):
        self._apply_env_and_defaults()
        self._ensure_directories()

    def _apply_env_and_defaults(self):
        for key, default in self.DEFAULTS.items():
            env_key = f'RAG_{key.upper()}'
            value = os.getenv(env_key, default)
            if isinstance(default, bool):
                value = str(value).lower() in ('true', '1', 'yes')
            elif isinstance(default, int):
                value = int(value)
            elif isinstance(default, float):
                value = float(value)
            setattr(self, key, value)

    def _ensure_directories(self):
        for folder in ['data_folder', 'cache_folder', 'backup_folder', 'log_folder']:
            Path(getattr(self, folder)).mkdir(parents=True, exist_ok=True)

    def load(self):
        config_file = Path(self.data_folder) / "config.json"
        if config_file.exists():
            try:
                with open(config_file, 'r', encoding='utf-8') as f:
                    data = json.load(f)
                    for k, v in data.items():
                        if hasattr(self, k):
                            setattr(self, k, v)
                logger.info("✅ Configuration loaded")
            except Exception as e:
                logger.error(f"⚠️ Failed to load config: {e}")

    def save(self):
        try:
            config_file = Path(self.data_folder) / "config.json"
            data = {k: v for k, v in self.__dict__.items() if k in self.DEFAULTS}
            with open(config_file, 'w', encoding='utf-8') as f:
                json.dump(data, f, indent=2)
            logger.info("✅ Configuration saved")
        except Exception as e:
            logger.error(f"⚠️ Failed to save config: {e}")

    def validate(self) -> Tuple[bool, List[str]]:
        issues = []
        try:
            resp = requests.get(f"{self.ollama_url}/api/tags", timeout=5)
            if resp.status_code != 200:
                issues.append(f"Ollama unreachable at {self.ollama_url}")
        except Exception as e:
            issues.append(f"Ollama connection failed: {e}")
        if not (100 <= self.chunk_size <= 2000):
            issues.append("chunk_size must be 100–2000")
        if not (1 <= self.top_k_retrieval <= 100):
            issues.append("top_k_retrieval must be 1–100")
        return len(issues) == 0, issues

config = MyConfig()

# 📝 LOGGING SETUP
class MyLogger:
    def __init__(self, log_dir: str = "./logs"):
        Path(log_dir).mkdir(exist_ok=True)
        log_file = Path(log_dir) / f"rag_{datetime.now().strftime('%Y%m%d')}.log"
        logging.basicConfig(
            level=logging.INFO,
            format='%(asctime)s | %(levelname)-8s | %(name)s | %(message)s',
            handlers=[
                logging.FileHandler(log_file, encoding='utf-8'),
                logging.StreamHandler(sys.stdout)
            ]
        )
        self.logger = logging.getLogger('RAG')

logger = MyLogger().logger

# 📚 DEPENDENCIES
try:
    import chromadb
    from chromadb.config import Settings
    from sentence_transformers import SentenceTransformer, CrossEncoder
    import requests
    import numpy as np
    from tqdm import tqdm
    from PyPDF2 import PdfReader
except ImportError as e:
    logger.critical(f"Missing core dependency: {e}")
    print("❌ Install: pip install chromadb sentence-transformers requests numpy tqdm PyPDF2")
    sys.exit(1)

# Optional format support
OPTIONAL_MODULES = {
    'docx': ('python-docx', False),
    'pptx': ('python-pptx', False),
    'openpyxl': ('openpyxl', False),
    'bs4': ('beautifulsoup4', False)
}
for mod, (pkg, _) in OPTIONAL_MODULES.items():
    try:
        __import__(mod)
        OPTIONAL_MODULES[mod] = (pkg, True)
    except ImportError:
        logger.info(f"Optional: {pkg} not installed → {mod.upper()} support disabled")

# 📖 DOCUMENT READER
class Reader:
    SUPPORTED_EXTENSIONS = {
        '.pdf', '.txt', '.md', '.csv', '.docx', '.doc', '.pptx', '.ppt',
        '.xlsx', '.xls', '.html', '.htm'
    }

    @staticmethod
    def read(file_path: str) -> Tuple[str, Dict]:
        ext = Path(file_path).suffix.lower()
        if ext not in Reader.SUPPORTED_EXTENSIONS:
            return "", {"error": "Unsupported file type"}
        try:
            if ext == '.pdf':
                text = Reader._read_pdf(file_path)
            elif ext in ('.docx', '.doc'):
                if not OPTIONAL_MODULES['docx'][1]:
                    return "", {"error": "Install python-docx"}
                from docx import Document as DocxDocument
                doc = DocxDocument(file_path)
                text = '\n'.join(p.text for p in doc.paragraphs if p.text.strip())
            elif ext in ('.pptx', '.ppt'):
                if not OPTIONAL_MODULES['pptx'][1]:
                    return "", {"error": "Install python-pptx"}
                from pptx import Presentation
                prs = Presentation(file_path)
                text = '\n'.join(
                    shape.text for slide in prs.slides for shape in slide.shapes
                    if hasattr(shape, "text") and shape.text.strip()
                )
            elif ext in ('.xlsx', '.xls'):
                if not OPTIONAL_MODULES['openpyxl'][1]:
                    return "", {"error": "Install openpyxl"}
                import openpyxl
                wb = openpyxl.load_workbook(file_path, data_only=True)
                text = '\n'.join(
                    f"[Sheet: {name}]\n" + '\n'.join(
                        ' | '.join(str(c) for c in row if c is not None)
                        for row in wb[name].iter_rows(values_only=True)
                    )
                    for name in wb.sheetnames
                )
            elif ext == '.csv':
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    reader = csv.reader(f)
                    text = '\n'.join(' | '.join(row) for row in reader)
            elif ext in ('.html', '.htm'):
                if not OPTIONAL_MODULES['bs4'][1]:
                    text = Reader._read_text(file_path)
                else:
                    from bs4 import BeautifulSoup
                    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                        soup = BeautifulSoup(f.read(), 'html.parser')
                        for tag in soup(["script", "style"]):
                            tag.decompose()
                        text = soup.get_text(separator='\n', strip=True)
            else:
                text = Reader._read_text(file_path)
            
            metadata = {
                'filename': Path(file_path).name,
                'extension': ext,
                'size': os.path.getsize(file_path),
                'char_count': len(text)
            }
            return text, metadata
        except Exception as e:
            logger.error(f"Read error {file_path}: {e}")
            return "", {"error": str(e)}

    @staticmethod
    def _read_text(path: str) -> str:
        for enc in ['utf-8', 'latin-1', 'cp1252']:
            try:
                with open(path, 'r', encoding=enc) as f:
                    return f.read()
            except UnicodeDecodeError:
                continue
        with open(path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()

    @staticmethod
    def _read_pdf(path: str) -> str:
        reader = PdfReader(path)
        return '\n'.join(
            f"[Page {i+1}]\n{page.extract_text()}"
            for i, page in enumerate(reader.pages)
            if page.extract_text()
        )

# ✂️ SMART CHUNKING
class SmartChunker:
    @staticmethod
    def chunk(text: str) -> List[Dict]:
        """Enhanced chunking with paragraph awareness"""
        text = re.sub(r'\s+', ' ', text).strip()
        if len(text) < config.chunk_size:
            return [{"text": text, "index": 0}]
        # Split on paragraph breaks first (stronger boundaries)
        paragraphs = [p.strip() for p in re.split(r'\n\s*\n', text) if p.strip()]
        chunks = []
        current_chunk = ""
        for para in paragraphs:
            sentences = re.split(r'(?<=[.!?])\s+', para)
            for sent in sentences:
                test = current_chunk + (' ' if current_chunk else '') + sent
                if len(test) > config.chunk_size and current_chunk:
                    # Save current chunk if it's substantial
                    if len(current_chunk.strip()) > 50:
                        chunks.append({"text": current_chunk.strip(), "index": len(chunks)})
                    # Create overlap with last few words
                    overlap_words = current_chunk.split()[-15:]
                    current_chunk = ' '.join(overlap_words) + ' ' + sent
                else:
                    current_chunk = test
            # Force boundary at paragraph end (preserve semantic units)
            if current_chunk.strip() and len(current_chunk.strip()) > 50:
                chunks.append({"text": current_chunk.strip(), "index": len(chunks)})
                # Keep small overlap for next paragraph
                overlap_words = current_chunk.split()[-10:]
                current_chunk = ' '.join(overlap_words)
        # Final chunk
        if current_chunk.strip() and len(current_chunk) > 50:
            chunks.append({"text": current_chunk.strip(), "index": len(chunks)})
        return chunks if chunks else [{"text": text, "index": 0}]

# 🧠 RAG ENGINE
class MyRAG:
    def __init__(self):
        logger.info("Initializing RAG System")
        self._validate_and_setup()
        self._load_models()
        self._init_database()
        self._load_registry_and_cache()
        self.history = []
        self.metrics = {'queries': 0, 'cache_hits': 0, 'errors': 0, 'avg_response_time': 0.0}
        self._warmup_model()
        logger.info("✅ System ready")

    def _validate_and_setup(self):
        ok, issues = config.validate()
        if not ok:
            for issue in issues:
                logger.error(f"❗ {issue}")
            sys.exit(1)
        for folder in [config.data_folder, config.cache_folder, config.backup_folder, config.log_folder]:
            Path(folder).mkdir(exist_ok=True)

    def _load_models(self):
        logger.info("🔥 Loading embedding model...")
        self.embed_model = SentenceTransformer(config.embedding_model)
        if config.use_reranking:
            logger.info("🔥 Loading reranker...")
            self.rerank_model = CrossEncoder("cross-encoder/ms-marco-MiniLM-L-6-v2")
        else:
            self.rerank_model = None

    def _init_database(self):
        self.db = chromadb.PersistentClient(
            path=config.data_folder,
            settings=Settings(anonymized_telemetry=False)
        )
        self.collection = self.db.get_or_create_collection(
            name="documents",
            metadata={"hnsw:space": "cosine"}
        )

    def _load_registry_and_cache(self):
        self.registry_file = Path(config.data_folder) / "doc_registry.json"
        self.cache_file = Path(config.cache_folder) / "answers.pkl"
        self.registry = self._load_json(self.registry_file)
        self.cache = self._load_pickle(self.cache_file) if config.enable_cache else {}
        if self.cache:
            logger.info(f"✅ Loaded {len(self.cache)} cached answers")

    def _load_json(self, path):
        if path.exists():
            try:
                with open(path, 'r', encoding='utf-8') as f:
                    return json.load(f)
            except Exception as e:
                logger.warning(f"Failed to load JSON: {e}")
        return {}

    def _load_pickle(self, path):
        if path.exists():
            try:
                with open(path, 'rb') as f:
                    return pickle.load(f)
            except Exception as e:
                logger.warning(f"Cache load failed: {e}")
        return {}

    def _save_json(self, path, data):
        with open(path, 'w', encoding='utf-8') as f:
            json.dump(data, f, indent=2)

    def _save_pickle(self, path, data):
        with open(path, 'wb') as f:
            pickle.dump(data, f)

    def _warmup_model(self):
        """Pre-load model to avoid cold start delays"""
        try:
            logger.info("🔥 Warming up model...")
            payload = {
                "model": config.model_name,
                "prompt": "Respond with 'OK'.",
                "stream": False,
                "options": {"temperature": 0, "num_predict": 5}
            }
            requests.post(f"{config.ollama_url}/api/generate", json=payload, timeout=10)
            logger.info("✅ Model ready")
        except:
            logger.warning("⚠️ Could not pre-warm model")

    def add_document(self, file_path: str, force: bool = False) -> bool:
        filename = Path(file_path).name
        # Check if already processed
        if not force and filename in self.registry:
            current_hash = self._file_hash(file_path)
            if current_hash == self.registry[filename].get('hash'):
                logger.info(f"📄 {filename} unchanged — skipping")
                return False
        text, meta = Reader.read(file_path)
        if not text.strip():
            logger.warning(f"📄 {filename} yielded no text")
            return False
        self._remove_document_chunks(filename)
        chunks = SmartChunker.chunk(text)
        logger.info(f"🧠 Embedding {len(chunks)} chunks for {filename}")
        embeddings = self.embed_model.encode(
            [c['text'] for c in chunks],
            batch_size=config.batch_size,
            normalize_embeddings=True,
            show_progress_bar=False
        ).tolist()
        timestamp = datetime.now().isoformat()
        ids = [f"{filename}_{timestamp}_{i}" for i in range(len(chunks))]
        metadatas = [{'source': filename, 'added': timestamp, 'chunk_idx': i} for i in range(len(chunks))]
        self.collection.add(
            ids=ids,
            documents=[c['text'] for c in chunks],
            embeddings=embeddings,
            metadatas=metadatas
        )
        self.registry[filename] = {
            'path': file_path,
            'chunks': len(chunks),
            'size': os.path.getsize(file_path),
            'added': timestamp,
            'hash': self._file_hash(file_path)
        }
        self._save_json(self.registry_file, self.registry)
        logger.info(f"✅ Added {filename}")
        return True

    def _file_hash(self, path: str) -> str:
        h = hashlib.md5()
        with open(path, 'rb') as f:
            for chunk in iter(lambda: f.read(8192), b""):
                h.update(chunk)
        return h.hexdigest()

    def _remove_document_chunks(self, filename: str):
        try:
            res = self.collection.get(where={"source": filename})
            if res['ids']:
                self.collection.delete(ids=res['ids'])
                logger.info(f"🗑️ Removed {len(res['ids'])} chunks for {filename}")
        except Exception as e:
            logger.error(f"❌ Failed to delete chunks: {e}")

    def ask(self, question: str, use_documents: bool = True, conversation_context: str = "") -> Dict:
        start = time.time()
        self.metrics['queries'] += 1
        # General knowledge mode (no RAG)
        if not use_documents:
            prompt = f"Answer naturally and conversationally:\n{question}"
            answer = self._generate(prompt, stream=config.stream_answers, mode='chat')
            return self._finalize_response(answer, [], start, mode='general')

        # Check cache
        cache_key = question.lower().strip()
        if config.enable_cache and cache_key in self.cache:
            cached = self.cache[cache_key]
            age = datetime.now() - datetime.fromisoformat(cached['time'])
            if age < timedelta(hours=config.cache_ttl_hours):
                self.metrics['cache_hits'] += 1
                if not config.stream_answers:
                    print("\n" + "*"*80)
                    print("📝 ANSWER (cached):\n")
                    print(cached['answer'])
                    print("*"*80)
                if config.show_sources:
                    print(f"\n💡 Sources: {', '.join(cached['sources'])}")
                return {
                    'answer': cached['answer'],
                    'sources': cached['sources'],
                    'time': time.time() - start,
                    'cached': True,
                    'mode': 'cached'
                }

        # Check if documents exist
        if self.collection.count() == 0:
            return self._finalize_response(
                "No documents in knowledge base. Add documents first.",
                [], start, mode='error'
            )

        # Enhance query with conversation context if provided
        search_query = question
        if conversation_context:
            search_query = f"{conversation_context}\nCurrent question: {question}"

        # Retrieve relevant documents
        logger.info("🔍 Searching documents...")
        emb = self.embed_model.encode(search_query, normalize_embeddings=True).tolist()
        results = self.collection.query(
            query_embeddings=[emb],
            n_results=config.top_k_retrieval,
            include=["documents", "metadatas"]
        )
        if not results['documents'][0]:
            return self._finalize_response("No relevant information found.", [], start, mode='no_results')

        docs = results['documents'][0]
        metas = results['metadatas'][0]
        # Rerank if enabled
        if self.rerank_model and config.use_reranking:
            logger.info("🎯 Reranking...")
            scores = self.rerank_model.predict([[question, d] for d in docs])
            idxs = np.argsort(scores)[::-1][:config.top_k_rerank]
            docs = [docs[i] for i in idxs]
            metas = [metas[i] for i in idxs]
        else:
            docs = docs[:config.top_k_rerank]
            metas = metas[:config.top_k_rerank]

        # Build context and generate answer
        context = "\n".join(f"[{m['source']}]\n{d}" for d, m in zip(docs, metas))
        prompt = self._build_prompt(question, context, conversation_context)
        answer = self._generate(prompt, stream=config.stream_answers, mode='rag')
        sources = list(set(m['source'] for m in metas))
        # Cache the result
        if config.enable_cache and answer and answer.strip():
            self.cache[cache_key] = {
                'answer': answer,
                'sources': sources,
                'time': datetime.now().isoformat()
            }
            self._save_pickle(self.cache_file, self.cache)
        return self._finalize_response(answer, sources, start, cached=False, mode='rag')

    def _build_prompt(self, question: str, context: str, conversation_context: str = "") -> str:
        base_instructions = """You are a helpful assistant that answers questions based on the provided context.
Instructions:
- Answer naturally and conversationally in complete sentences
- Use information ONLY from the context provided
- Be concise but complete - include all relevant details
- If asked about people, actions, or events, be specific and detailed
- When the word 'with' appears, consider context: people involved, manner/method, or tools used
- If information is insufficient, say so clearly
- Cite sources naturally, e.g., "According to [filename]..."
- Never say "the context doesn't mention" - instead use "Based on available information..."
"""
        if conversation_context:
            return f"""{base_instructions}
Previous conversation context:
{conversation_context}
Document Context:
{context}
Current Question: {question}
Answer:"""
        else:
            return f"""{base_instructions}
Context:
{context}
Question: {question}
Answer:"""

    def _generate(self, prompt: str, stream: bool = False, mode: str = 'rag') -> str:
        try:
            url = f"{config.ollama_url}/api/generate"
            payload = {
                "model": config.model_name,
                "prompt": prompt,
                "stream": stream,
                "options": {
                    "temperature": 0.7 if mode == 'chat' else 0.5,
                    "num_predict": 512,
                    "num_ctx": 3072,
                    "repeat_penalty": 1.1,
                    "top_k": 40,
                    "top_p": 0.9
                }
            }
            if not stream:
                resp = requests.post(url, json=payload, timeout=config.ollama_timeout)
                if resp.status_code == 200:
                    ans = resp.json().get('response', '').strip()
                    return ans if ans else "I couldn't generate a clear answer."
                else:
                    logger.error(f"Ollama error {resp.status_code}: {resp.text}")
                    return f"❌ Generation failed (status {resp.status_code})"
            else:
                full = ""
                try:
                    with requests.post(url, json=payload, stream=True, timeout=config.ollama_timeout) as r:
                        if r.status_code != 200:
                            logger.error(f"Streaming failed: {r.status_code}")
                            return f"❌ Streaming error: {r.status_code}"
                        print("\n📝 ANSWER:\n")
                        for line in r.iter_lines():
                            if line:
                                try:
                                    chunk = json_lib.loads(line.decode())
                                    token = chunk.get("response", "")
                                    if token:
                                        print(token, end="", flush=True)
                                        full += token
                                    if chunk.get("done"):
                                        break
                                except Exception as je:
                                    logger.debug(f"Parse error: {je}")
                                    continue
                        print("\n" + "*"*80)
                except Exception as e:
                    logger.error(f"Stream error: {e}")
                    return "❌ Connection interrupted during generation."
                return full.strip() if full.strip() else "I retrieved data but couldn't generate a proper answer."
        except Exception as e:
            logger.error(f"Generate error: {e}", exc_info=True)
            self.metrics['errors'] += 1
            return "❌ An internal error occurred while generating the answer."

    def _finalize_response(self, answer, sources, start_time, cached=False, mode='rag'):
        elapsed = time.time() - start_time
        # Ensure answer is never empty
        if not answer or not answer.strip():
            answer = "I found relevant information but couldn't generate a clear answer."
        # Display answer (if not already streamed)
        if not cached and not config.stream_answers:
            print("\n" + "*"*80)
            print("📝 ANSWER:\n")
            print(answer)
            print("*"*80)
        # Show sources
        if config.show_sources and sources:
            print(f"\n💡 Sources: {', '.join(sources)}")
        # Show timing
        if not cached:
            print(f"\n⏱️ Response time: {elapsed:.2f}s")
            # Update average response time
            avg = self.metrics['avg_response_time']
            n = self.metrics['queries']
            self.metrics['avg_response_time'] = (avg * (n - 1) + elapsed) / n
        return {
            'answer': answer,
            'sources': sources,
            'time': elapsed,
            'cached': cached,
            'mode': mode
        }

    def list_documents(self):
        docs = []
        for name, info in self.registry.items():
            added = info.get('added', datetime.now().isoformat())
            docs.append({
                'filename': name,
                'chunks': info.get('chunks', 0),
                'size_mb': info.get('size', 0) / (1024**2),
                'added': added
            })
        return docs

    def remove_document(self, filename: str) -> bool:
        if filename not in self.registry:
            return False
        self._remove_document_chunks(filename)
        del self.registry[filename]
        self._save_json(self.registry_file, self.registry)
        # Invalidate related cache entries
        if config.enable_cache:
            to_del = [k for k, v in self.cache.items() if filename in v.get('sources', [])]
            for k in to_del:
                del self.cache[k]
            if to_del:
                self._save_pickle(self.cache_file, self.cache)
                logger.info(f"Invalidated {len(to_del)} cache entries")
        logger.info(f"✅ Removed {filename}")
        return True

    def get_stats(self):
        cache_rate = 0
        if self.metrics['queries']:
            cache_rate = (self.metrics['cache_hits'] / self.metrics['queries']) * 100
        return {
            'documents': len(self.registry),
            'chunks': self.collection.count(),
            'cache_entries': len(self.cache),
            'queries': self.metrics['queries'],
            'cache_hit_rate': f"{cache_rate:.1f}%",
            'avg_response_time': f"{self.metrics['avg_response_time']:.2f}s",
            'errors': self.metrics['errors']
        }

    def clear_cache(self):
        count = len(self.cache)
        self.cache = {}
        if config.enable_cache:
            self._save_pickle(self.cache_file, self.cache)
        logger.info(f"🧹 Cleared {count} cache entries")
        return count

    def backup(self):
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        backup_dir = Path(config.backup_folder) / f"backup_{timestamp}"
        backup_dir.mkdir(parents=True)
        shutil.copytree(config.data_folder, backup_dir / "data", dirs_exist_ok=True)
        if config.enable_cache and self.cache_file.exists():
            shutil.copy2(self.cache_file, backup_dir / "cache.pkl")
        info = {
            'timestamp': timestamp,
            'documents': len(self.registry),
            'chunks': self.collection.count(),
            'cache_entries': len(self.cache)
        }
        with open(backup_dir / "info.json", 'w') as f:
            json.dump(info, f, indent=2)
        logger.info(f"✅ Backup saved to {backup_dir}")
        return str(backup_dir)

# 💬 CLI
class CLI:
    def __init__(self):
        self.rag = None
        self.running = True
        signal.signal(signal.SIGINT, self._exit_gracefully)
        signal.signal(signal.SIGTERM, self._exit_gracefully)

    def _exit_gracefully(self, signum, frame):
        print("\n👋 Shutting down gracefully...")
        if self.rag and config.auto_backup:
            print("Creating automatic backup...")
            try:
                self.rag.backup()
            except:
                pass
        sys.exit(0)

    def run(self):
        self._print_banner()
        config.load()
        ok, issues = config.validate()
        if not ok:
            print("\n⚠️  Configuration issues:")
            for issue in issues:
                print(f" • {issue}")
            if not self._confirm("Continue anyway?"):
                return
        try:
            self.rag = MyRAG()
        except Exception as e:
            logger.critical(f"Init failed: {e}")
            print(f"❌ Fatal: {e}")
            sys.exit(1)
        while self.running:
            self._main_menu()

    def _print_banner(self):
        print("\n" + "="*80)
        print("🚀 RAG SYSTEM ")
        print("="*80)
        print(f"📁 Data: {config.data_folder}")
        print(f"🧠 Model: {config.model_name}")
        print(f"🔍 Embedding: {config.embedding_model}")
        print("="*80)

    def _main_menu(self):
        print("\n" + "*"*80)
        print("📋 MAIN MENU")
        print("*"*80)
        print("1. 💬 Chat Mode")
        print("2. ➕ Add Document")
        print("3. 📁 Add Folder")
        print("4. 📚 List Documents")
        print("5. 🗑️  Remove Document")
        print("6. 📊 Stats")
        print("7. 🏥 Health Check")
        print("8. 💾 Backup")
        print("9. 🧹 Clear Cache")
        print("10. ⚙️ Settings")
        print("0. 🚪 Exit")
        choice = input("\n👉 Choose (0-10): ").strip()
        actions = {
            '1': self._chat_mode,
            '2': self._add_document,
            '3': self._add_folder,
            '4': self._list_docs,
            '5': self._remove_doc,
            '6': self._show_stats,
            '7': self._health_check,
            '8': self._backup,
            '9': self._clear_cache,
            '10': self._settings,
            '0': self._quit
        }
        actions.get(choice, lambda: print("❌ Invalid choice"))()

    def _chat_mode(self):
        print("\n" + "*"*80)
        print("💬 CHAT MODE")
        print("*"*80)
        print("Commands:")
        print("  /help    - Show tips")
        print("  /docs    - List documents")
        print("  /stats   - System statistics")
        print("  /clear   - Clear conversation memory")
        print("  /exit    - Return to main menu")
        print("*"*80)
        # Conversation history for context
        conversation_history = []
        while True:
            q = input("\nYou: ").strip()
            if not q:
                continue
            # Handle commands
            cmd_lower = q.lower()
            if cmd_lower in ['exit', 'quit', '/exit', '/quit']:
                return
            if cmd_lower in ['back', '/back']:
                break
            if cmd_lower in ['stats', '/stats']:
                self._show_stats()
                continue
            if cmd_lower in ['docs', '/docs']:
                self._list_docs()
                continue
            if cmd_lower in ['clear', '/clear']:
                conversation_history.clear()
                print("✅ Conversation memory cleared")
                continue
            if cmd_lower in ['help', '/help']:
                print("\n💡 Tips:")
                print(" • Ask questions about your documents")
                print(" • Use 'Y' to search documents, 'n' for general knowledge")
                print(" • The system remembers conversation context")
                print(" • Use /clear to reset conversation memory")
                continue
            # Add user message to history
            conversation_history.append({"role": "user", "content": q})
            # Determine if we should use documents
            use_documents = False
            if self.rag.collection.count() > 0:
                use_docs_input = input("Use document knowledge? (Y/n, default=Y): ").strip().lower()
                use_documents = use_docs_input not in ('n', 'no')
            else:
                print("📭 No documents available - using general knowledge")
            print()
            # Generate response
            if not use_documents:
                # Use chat API with conversation history
                try:
                    system_prompt = """You are a helpful, friendly, and conversational assistant. Respond in complete, natural sentences. Remember context from the conversation and reference it when relevant. Be warm, concise, and human-like."""
                    messages_with_system = [{"role": "system", "content": system_prompt}] + conversation_history[-10:]
                    payload = {
                        "model": config.model_name,
                        "messages": messages_with_system,
                        "stream": config.stream_answers,
                        "options": {
                            "temperature": 0.7,
                            "num_ctx": 4096,
                            "repeat_penalty": 1.1
                        }
                    }
                    if not config.stream_answers:
                        resp = requests.post(f"{config.ollama_url}/api/chat", json=payload, timeout=config.ollama_timeout)
                        if resp.status_code == 200:
                            answer = resp.json().get('message', {}).get('content', '').strip()
                            answer = answer if answer else "I couldn't generate a response."
                            print(f"\n{answer}")
                            conversation_history.append({"role": "assistant", "content": answer})
                        else:
                            error_msg = f"❌ Error: {resp.status_code}"
                            print(error_msg)
                            conversation_history.append({"role": "assistant", "content": error_msg})
                    else:
                        full = ""
                        print("\n📝 Response:\n")
                        with requests.post(f"{config.ollama_url}/api/chat", json=payload, stream=True, timeout=config.ollama_timeout) as r:
                            if r.status_code != 200:
                                error_msg = f"❌ Streaming error: {r.status_code}"
                                print(error_msg)
                                conversation_history.append({"role": "assistant", "content": error_msg})
                                continue
                            for line in r.iter_lines():
                                if line:
                                    try:
                                        chunk = json_lib.loads(line.decode())
                                        token = chunk.get("message", {}).get("content", "")
                                        if token:
                                            print(token, end="", flush=True)
                                            full += token
                                        if chunk.get("done"):
                                            break
                                    except:
                                        continue
                        print("\n" + "*"*80)
                        full = full.strip() if full.strip() else "I couldn't generate a response."
                        conversation_history.append({"role": "assistant", "content": full})
                except Exception as e:
                    error_msg = f"❌ Generation error: {e}"
                    print(error_msg)
                    conversation_history.append({"role": "assistant", "content": error_msg})
            else:
                # Use RAG with conversation context
                context_str = "\n".join([
                    f"{m['role'].title()}: {m['content']}" 
                    for m in conversation_history[-6:]
                ])
                result = self.rag.ask(q, use_documents=True, conversation_context=context_str)
                conversation_history.append({"role": "assistant", "content": result['answer']})

    def _add_document(self):
        path = input("\n📄 File path: ").strip().strip('"')
        if not path:
            print("❌ No path provided")
            return
        if not Path(path).exists():
            print("❌ File not found")
            return
        force = self._confirm("Force reprocess (ignore cache)?")
        try:
            if self.rag.add_document(path, force=force):
                print("✅ Document added successfully")
            else:
                print("ℹ️  Skipped (unchanged or empty)")
        except Exception as e:
            print(f"❌ Error adding document: {e}")

    def _add_folder(self):
        path = input("\n📁 Folder path: ").strip().strip('"')
        if not path:
            print("❌ No path provided")
            return
        if not Path(path).exists():
            print("❌ Folder not found")
            return
        recursive = self._confirm("Include subfolders?")
        pattern = "**/*" if recursive else "*"
        files = [
            f for f in Path(path).glob(pattern) 
            if f.is_file() and f.suffix.lower() in Reader.SUPPORTED_EXTENSIONS
        ]
        if not files:
            print("⚠️ No supported files found")
            return
        print(f"\n📚 Found {len(files)} files")
        if not self._confirm("Process all files?"):
            return
        print(f"\n🔄 Processing {len(files)} files...")
        success = 0
        skipped = 0
        errors = 0
        for f in tqdm(files, desc="Adding documents"):
            try:
                if self.rag.add_document(str(f)):
                    success += 1
                else:
                    skipped += 1
            except Exception as e:
                errors += 1
                logger.error(f"Failed to add {f}: {e}")
        print(f"\n✅ Folder processed:")
        print(f"   • Added: {success}")
        print(f"   • Skipped: {skipped}")
        print(f"   • Errors: {errors}")

    def _list_docs(self):
        docs = self.rag.list_documents()
        if not docs:
            print("\n📭 No documents in knowledge base")
            return
        print("\n" + "*"*80)
        print("📚 DOCUMENTS")
        print("*"*80)
        # Sort by date 
        sorted_docs = sorted(docs, key=lambda x: x['added'], reverse=True)
        for i, d in enumerate(sorted_docs, 1):
            print(f"\n{i}. {d['filename']}")
            print(f"   📄 {d['chunks']} chunks | 💾 {d['size_mb']:.2f} MB | 🕐 {d['added'][:19]}")
        print(f"\n📊 Total: {len(docs)} documents, {sum(d['chunks'] for d in docs)} chunks")

    def _remove_doc(self):
        docs = self.rag.list_documents()
        if not docs:
            print("\n📭 No documents to remove")
            return
        self._list_docs()
        name = input("\n🗑️ Filename to remove: ").strip()
        if name not in [d['filename'] for d in docs]:
            print("❌ Document not found")
            return
        if self._confirm(f"Remove '{name}'?"):
            if self.rag.remove_document(name):
                print("✅ Document removed")
            else:
                print("❌ Failed to remove document")

    def _show_stats(self):
        stats = self.rag.get_stats()
        print("\n" + "*"*80)
        print("📊 SYSTEM STATISTICS")
        print("*"*80)
        for k, v in stats.items():
            label = k.replace('_', ' ').title()
            print(f"{label}: {v}")
        # Additional info
        print(f"\nModel: {config.model_name}")
        print(f"Chunk Size: {config.chunk_size} (overlap: {config.chunk_overlap})")
        print(f"Top-K Retrieval: {config.top_k_retrieval} → Rerank: {config.top_k_rerank}")

    def _health_check(self):
        print("\n🏥 System Health Check")
        print("*"*80)
        # Check Ollama
        try:
            resp = requests.get(f"{config.ollama_url}/api/tags", timeout=5)
            ollama_ok = resp.status_code == 200
            if ollama_ok:
                models = resp.json().get('models', [])
                model_exists = any(m.get('name') == config.model_name for m in models)
                if not model_exists:
                    print(f"⚠️  Ollama: Connected, but model '{config.model_name}' not found")
                    ollama_ok = False
        except Exception as e:
            ollama_ok = False
            print(f"❌ Ollama: {e}")
        # Check Database
        try:
            db_ok = self.rag.collection.count() >= 0
        except:
            db_ok = False
        # Check Disk Space
        try:
            usage = shutil.disk_usage(config.data_folder)
            disk_ok = usage.free > 100 * 1024 * 1024 
            free_gb = usage.free / (1024**3)
        except:
            disk_ok = False
            free_gb = 0
        # Status
        all_ok = all([ollama_ok, db_ok, disk_ok])
        status = "✅ HEALTHY" if all_ok else "⚠️ DEGRADED"
        print(f"\nStatus: {status}")
        print(f" • Ollama: {'✅' if ollama_ok else '❌'} ({config.ollama_url})")
        print(f" • Database: {'✅' if db_ok else '❌'} ({self.rag.collection.count()} chunks)")
        print(f" • Disk space: {'✅' if disk_ok else '❌'} ({free_gb:.1f} GB free)")
        print(f" • Cache: {'✅ Enabled' if config.enable_cache else '⚠️ Disabled'}")
        print(f" • Reranking: {'✅ Enabled' if config.use_reranking else '⚠️ Disabled'}")

    def _backup(self):
        if not self._confirm("Create backup now?"):
            return
        try:
            print("\n💾 Creating backup...")
            path = self.rag.backup()
            print(f"✅ Backup saved to: {path}")
        except Exception as e:
            print(f"❌ Backup failed: {e}")

    def _clear_cache(self):
        count = len(self.rag.cache)
        if count == 0:
            print("\nℹ️  Cache is already empty")
            return
        if self._confirm(f"Clear {count} cached answers?"):
            cleared = self.rag.clear_cache()
            print(f"✅ Cleared {cleared} cache entries")

    def _settings(self):
        print("\n" + "*"*80)
        print("⚙️ SETTINGS")
        print("*"*80)
        settings = [
            ("Show sources", "show_sources"),
            ("Stream answers", "stream_answers"),
            ("Use reranking", "use_reranking"),
            ("Enable cache", "enable_cache"),
            ("Auto backup", "auto_backup"),
            ("Model name", "model_name"),
            ("Chunk size", "chunk_size"),
            ("Chunk overlap", "chunk_overlap"),
            ("Top-K retrieval", "top_k_retrieval"),
            ("Top-K rerank", "top_k_rerank"),
            ("Timeout (s)", "ollama_timeout")
        ]
        for i, (label, attr) in enumerate(settings, 1):
            val = getattr(config, attr)
            if isinstance(val, bool):
                val = "ON" if val else "OFF"
            print(f"{i:2}. {label:20} : {val}")
        choice = input("\n👉 Edit setting # (or Enter to skip): ").strip()
        if not choice.isdigit() or not (1 <= int(choice) <= len(settings)):
            return
        label, attr = settings[int(choice) - 1]
        current = getattr(config, attr)
        if isinstance(current, bool):
            setattr(config, attr, not current)
            new_val = "ON" if getattr(config, attr) else "OFF"
            print(f"✅ {label} toggled to {new_val}")
        else:
            new = input(f"New value for {label} [{current}]: ").strip()
            if new:
                try:
                    if isinstance(current, int):
                        setattr(config, attr, int(new))
                    elif isinstance(current, float):
                        setattr(config, attr, float(new))
                    else:
                        setattr(config, attr, new)
                    print(f"✅ {label} updated to {new}")
                except ValueError:
                    print("❌ Invalid value")
                    return
        config.save()

    def _confirm(self, msg: str) -> bool:
        response = input(f"{msg} (y/N): ").strip().lower()
        return response in ('y', 'yes')

    def _quit(self):
        print("\n👋 Shutting down...")
        # Show final stats
        stats = self.rag.get_stats()
        print(f"Session summary: {stats['queries']} queries, {stats['cache_hit_rate']} cache hits")
        # Auto backup if enabled
        if config.auto_backup and stats['queries'] > 0:
            if self._confirm("Create backup before exit?"):
                try:
                    self.rag.backup()
                except:
                    pass
        self.running = False
        print("Goodbye! 👋")

# ▶️ MAIN
def main():
    try:
        cli = CLI()
        cli.run()
    except KeyboardInterrupt:
        print("\n👋 Interrupted by user")
    except Exception as e:
        logger.critical(f"Fatal error: {e}", exc_info=True)
        print(f"\n💥 Fatal error: {e}")
        print(f"Check logs in {config.log_folder}")
        sys.exit(1)

if __name__ == "__main__":
    main()
